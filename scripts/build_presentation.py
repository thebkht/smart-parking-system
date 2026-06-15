#!/usr/bin/env python3
"""Build the final Smart Parking System presentation (.pptx).

Styled to match the technical report (ICCV-style academic paper): Times serif,
black text on white, bold numbered section headings, booktabs-style tables, and
"Figure N." captions — no colored cards.

Scope: this is the *final* talk that follows the earlier progress review, so it
deliberately excludes topics already presented there (architecture, ML results,
evaluation depth, edge pipeline, backend endpoints, benchmarks, app screens) and
focuses on the new integration work plus the report's synthesis.

    .venv/bin/python scripts/build_presentation.py
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "outputs" / "figures"
OUT = ROOT / "outputs" / "Smart-Parking-System-Final-Presentation.pptx"

# ---- report palette: black serif on white ------------------------------
BLACK = RGBColor(0x00, 0x00, 0x00)
INK = RGBColor(0x11, 0x11, 0x11)
GRAY = RGBColor(0x55, 0x55, 0x55)
LINK = RGBColor(0x1A, 0x4E, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SERIF = "Times New Roman"
MONO = "Courier New"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5
ML = 0.85  # left margin


def slide():
    return prs.slides.add_slide(BLANK)


def _set(p, text, size, color=BLACK, bold=False, italic=False, font=SERIF, align=None):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    if align is not None:
        p.alignment = align
    return r


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    return tf


def rule(s, x, y, w, thick=0.018, color=BLACK):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(thick))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def heading(s, num, text):
    tf = textbox(s, ML, 0.52, SW - 2 * ML, 0.7)
    _set(tf.paragraphs[0], f"{num}.  {text}", 26, BLACK, bold=True)
    rule(s, ML, 1.28, SW - 2 * ML, 0.02)


def footer(s, n):
    tf = textbox(s, ML, 7.06, SW - 2 * ML, 0.3)
    _set(tf.paragraphs[0], "Smart Parking System — Final Technical Report", 9, GRAY, italic=True)
    tf2 = textbox(s, SW - ML - 0.8, 7.06, 0.8, 0.3)
    _set(tf2.paragraphs[0], str(n), 9, GRAY, align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, h, items, size=15, gap=10):
    """items: (head, body) — body may be None; or (None, sub) for a sub-bullet."""
    tf = textbox(s, x, y, w, h)
    first = True
    for head, body in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        if head is None:  # sub-bullet
            _set(p, "      – ", size - 1, GRAY)
            _set(p, body, size - 2, GRAY)
            continue
        _set(p, "•  ", size, BLACK, bold=True)
        _set(p, head, size, INK, bold=bool(body))
        if body:
            _set(p, "  " + body, size - 1.5, INK)
    return tf


def image(s, path, x, y, w, h, caption=None):
    iw, ih = Image.open(path).size
    ar = iw / ih
    if ar > w / h:
        nw, nh = w, w / ar
    else:
        nh, nw = h, h * ar
    px, py = x + (w - nw) / 2, y + (h - nh) / 2
    s.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(nw))
    if caption:
        tf = textbox(s, x, y + h + 0.03, w, 0.4)
        _set(tf.paragraphs[0], caption, 11, GRAY, italic=True, align=PP_ALIGN.CENTER)


def cell(s, x, y, w, h, text, size=12, bold=False, align=PP_ALIGN.CENTER, mono=False):
    tf = textbox(s, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    _set(tf.paragraphs[0], text, size, BLACK, bold=bold, font=MONO if mono else SERIF, align=align)


def booktabs(s, x, y, headers, rows, colw, row_h=0.46, header_h=0.46):
    xs, cx = [], x
    for cw in colw:
        xs.append(cx)
        cx += cw
    total = sum(colw)
    rule(s, x, y, total, 0.026)                       # top rule (thick)
    for j, htext in enumerate(headers):
        cell(s, xs[j], y + 0.05, colw[j], header_h, htext, 13, bold=True,
             align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    mid = y + header_h + 0.06
    rule(s, x, mid, total, 0.012)                     # mid rule (thin)
    ry = mid + 0.04
    for row in rows:
        for j, v in enumerate(row):
            cell(s, xs[j], ry, colw[j], row_h, v.replace("**", ""), 12.5,
                 bold="**" in v, align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
        ry += row_h
    rule(s, x, ry + 0.02, total, 0.026)               # bottom rule (thick)
    return ry


def codebox(s, x, y, w, h, lines, header=None):
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.shadow.inherit = False
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = BLACK
    box.line.width = Pt(0.75)
    tf = textbox(s, x + 0.18, y + 0.12, w - 0.36, h - 0.24)
    start = True
    if header:
        _set(tf.paragraphs[0], header, 11, GRAY, italic=True)
        start = False
    for ln in lines:
        p = tf.paragraphs[0] if start else tf.add_paragraph()
        start = False
        _set(p, ln, 11.5, BLACK, font=MONO)


# =========================================================================
# 1. Title (paper title block)
# =========================================================================
s = slide()
tf = textbox(s, 1.0, 1.9, SW - 2.0, 2.0, anchor=MSO_ANCHOR.TOP)
_set(tf.paragraphs[0], "Smart Parking System — Final Technical Report", 34, BLACK, bold=True, align=PP_ALIGN.CENTER)
p = tf.add_paragraph(); p.space_before = Pt(10)
_set(p, "Edge Occupancy Detection with Quadrilateral Pooling,", 19, INK, align=PP_ALIGN.CENTER)
p = tf.add_paragraph()
_set(p, "YOLOv8-cls, and a Find My Car Application", 19, INK, align=PP_ALIGN.CENTER)
rule(s, 4.67, 4.35, 4.0, 0.018)
auth = textbox(s, 1.0, 4.6, SW - 2.0, 1.2)
_set(auth.paragraphs[0], "Bakhtiyor Ganijon  ·  Sadriddinov Otabek  ·  Sattor Mamatov  ·  Mirzaev Komronkhon", 15, INK, align=PP_ALIGN.CENTER)
p = auth.add_paragraph(); p.space_before = Pt(6)
_set(p, "Final submission — follow-up to the progress review", 13, GRAY, italic=True, align=PP_ALIGN.CENTER)
p = auth.add_paragraph(); p.space_before = Pt(14)
_set(p, "github.com/thebkht/smart-parking-system", 13, LINK, font=MONO, align=PP_ALIGN.CENTER)

# =========================================================================
# 2. Problem statement (report §2)
# =========================================================================
s = slide()
heading(s, 2, "Problem statement")
bullets(s, ML, 1.7, SW - 2 * ML, 4.8, [
    ("Detect, per parking space, whether it is occupied or free —",
     "from a single overhead camera, without per-space in-ground sensors."),
    ("Keep imagery on the device.",
     "Streaming video to the cloud carries bandwidth and privacy cost; only "
     "compact occupancy JSON should leave the edge."),
    ("Generalize to unseen lots.",
     "Performance must hold on cameras and layouts never seen in training, not "
     "just on memorized scenes."),
    ("Close the product loop.",
     "Configuration (owner setup), live maps, and a photo-based Find My Car must "
     "all run off one consistent layout."),
])
footer(s, 2)

# =========================================================================
# 3. Since the progress review (scope of this talk)
# =========================================================================
s = slide()
heading(s, 3, "Since the progress review")
tf = textbox(s, ML, 1.55, SW - 2 * ML, 0.5)
_set(tf.paragraphs[0],
     "The model and benchmarks were presented earlier. This talk covers the integration work that "
     "completed the product, and the report's synthesis.", 14, GRAY, italic=True)
bullets(s, ML, 2.45, SW - 2 * ML, 3.6, [
    ("Owner setup now runs Structure-from-Motion server-side", "with a manual fallback (§4)."),
    ("Owner label correction", "— rename spots after layout generation (§4)."),
    ("Find My Car reference-photo management", "— per-spot references in the backend (§5)."),
    ("Optional authentication and session ownership", "(§6)."),
    ("A representative integration bug, and the report's Discussion + Conclusion", "(§7–§9)."),
], size=15, gap=12)
footer(s, 3)

# =========================================================================
# 4. Owner setup — server-side SfM + correction
# =========================================================================
s = slide()
heading(s, 4, "Owner setup — server-side SfM and correction")
bullets(s, ML, 1.6, 6.4, 4.8, [
    ("Upload and reconstruct",
     "POST /layout runs the SfM pass in-process to build a bird's-eye map and "
     "extract quadrilateral spot polygons, then persists the layout."),
    ("Graceful fallback",
     "If SfM cannot build a usable layout it returns 422; the app falls back to "
     "manual polygon submission via POST /map."),
    ("Inline label correction",
     "PATCH /spots/{id} lets the owner rename each spot after generation; the new "
     "label is reflected on the map and reused by Find My Car."),
])
image(s, FIG / "ui_owner_setup_layout.png", 7.55, 1.6, 5.0, 4.6,
      "Figure 1. Owner setup: generated layout ready to publish.")
footer(s, 4)

# =========================================================================
# 5. Find My Car — reference-photo management
# =========================================================================
s = slide()
heading(s, 5, "Find My Car — reference-photo management")
bullets(s, ML, 1.6, 6.6, 3.6, [
    ("Per-spot references in the backend",
     "POST /spots/{id}/references stores photos under a managed directory and "
     "records them; GET lists them."),
    ("Localization uses the managed set",
     "POST /park matches the driver photo against the managed references, falling "
     "back to the bundled sample set only when empty."),
    ("Unchanged localizer",
     "SIFT + FLANN + RANSAC, training-free and CPU-only."),
])
tf = textbox(s, ML, 5.4, 6.6, 1.4)
_set(tf.paragraphs[0], "Localizer accuracy (labelled night-overhead set):", 13, INK, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(4)
_set(p, "21/21 top-1  ·  21/21 top-3  ·  536 ms average  ·  no training", 13, INK)
image(s, FIG / "ui_find_my_car.png", 8.2, 1.55, 4.3, 5.0,
      "Figure 2. Find My Car: matched spot highlighted in amber.")
footer(s, 5)

# =========================================================================
# 6. Auth & session ownership
# =========================================================================
s = slide()
heading(s, 6, "Authentication and session ownership")
bullets(s, ML, 1.6, 6.5, 4.4, [
    ("Opt-in by design",
     "Off by default so the demo runs token-free; AUTH_ENABLED turns it on."),
    ("Owner-mutating routes are protected",
     "POST /map, POST /layout, PATCH /spots/{id}, and reference upload require a "
     "bearer token from POST /auth/register."),
    ("Sessions are owned",
     "POST /park stamps the owner; GET /find/{id} returns 404 for sessions owned "
     "by another user. Read-only /status and /map stay public."),
])
codebox(s, 7.4, 1.7, 5.1, 3.3,
        ['POST /auth/register',
         '  {"username": "owner"}',
         '  -> {"token": "..."}',
         '',
         'Authorization: Bearer <token>',
         '  on POST /map, PATCH /spots/{id},',
         '     POST /spots/{id}/references',
         '',
         'AUTH_ENABLED=1  (backend env)'],
        header="Lightweight token flow")
footer(s, 6)

# =========================================================================
# 7. Case study — live occupancy counter bug (report §6.4)
# =========================================================================
s = slide()
heading(s, 7, "Case study — the live-occupancy counter bug")
bullets(s, ML, 1.7, SW - 2 * ML, 4.6, [
    ("Symptom",
     "The web dashboard showed 55 free and 127 occupied against a 12-spot layout."),
    ("Root cause",
     "The counters aggregated every entry in the raw /status payload — including "
     "stale spot IDs from earlier edge runs — while the map rendered only the "
     "published layout's spots."),
    ("Fix",
     "Constrain the counters to the layout's spot IDs, restoring the invariant "
     "free + occupied ≤ total and matching the (already-correct) mobile client."),
    ("Hardening",
     "Recommend the backend filter /status responses to the published layout so "
     "all clients are consistent by construction."),
])
footer(s, 7)

# =========================================================================
# 8. Discussion (report §5)
# =========================================================================
s = slide()
heading(s, 8, "Discussion")
cols = [
    ("What worked", [
        "YOLOv8n beats the s/m variants on test",
        "Quad warp beats square crop by +1.34 pp",
        "Matches ResNet50 at ~9× fewer parameters",
    ]),
    ("Limitations", [
        "Patch quality is the ceiling (FN 26/605)",
        "Low-light recall drops to 93.5%",
        "Localization set is night-only, n = 21",
    ]),
    ("Production", [
        "858 FPS — ~300× reporting headroom",
        "99.9% bandwidth saving; 30-min stable",
        "SIFT-only Find My Car is the final scope",
    ]),
]
for i, (head, items) in enumerate(cols):
    x = ML + i * 4.05
    tfh = textbox(s, x, 1.7, 3.8, 0.5)
    _set(tfh.paragraphs[0], head, 17, BLACK, bold=True, italic=True)
    rule(s, x, 2.2, 3.6, 0.014)
    tf = textbox(s, x, 2.4, 3.85, 3.6)
    for j, it in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        _set(p, "•  ", 13, BLACK, bold=True)
        _set(p, it, 13.5, INK)
footer(s, 8)

# =========================================================================
# 9. Conclusion & future work (report §7)
# =========================================================================
s = slide()
heading(s, 9, "Conclusion and future work")
bullets(s, ML, 1.7, SW - 2 * ML, 3.6, [
    ("Delivered",
     "An edge system that lifts a 90.9% fixed-ROI baseline to 97.72% on the "
     "unseen-lot test split via quadrilateral pooling."),
    ("A complete product",
     "Owner setup with server-side SfM, web + mobile live maps, and photo-based "
     "Find My Car — validated by an in-process smoke test and client contract tests."),
    ("Future work",
     "Low-light augmentation, higher-fidelity SfM reconstruction, and a longer "
     "soak test with backend-side status filtering."),
])
rule(s, ML, 5.45, SW - 2 * ML, 0.02)
tf = textbox(s, ML, 5.65, SW - 2 * ML, 0.9)
_set(tf.paragraphs[0], "Thank you.", 18, BLACK, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(6)
_set(p, "Code: ", 14, INK)
_set(p, "github.com/thebkht/smart-parking-system", 14, LINK, font=MONO)
footer(s, 9)

prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")
